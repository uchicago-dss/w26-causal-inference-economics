from __future__ import annotations

import argparse
from pathlib import Path

import pypdfium2 as pdfium
from pptx import Presentation
from pptx.util import Inches


def export_pdf_to_pptx(pdf_path: Path, pptx_path: Path, dpi: int = 200) -> None:
    pdf_path = pdf_path.resolve()
    pptx_path = pptx_path.resolve()
    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF not found: {pdf_path}")

    pdf = pdfium.PdfDocument(str(pdf_path))

    prs = Presentation()
    # 16:9 widescreen to match your Beamer aspectratio=169.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for page_index in range(len(pdf)):
        page = pdf.get_page(page_index)
        bitmap = page.render(scale=dpi / 72.0)
        pil_image = bitmap.to_pil()

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(_save_temp_png(pil_image, pptx_path.parent, page_index)),
            left=Inches(0),
            top=Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        page.close()

    prs.save(str(pptx_path))


def _save_temp_png(pil_image, out_dir: Path, page_index: int) -> Path:
    out_dir.mkdir(parents=True, exist_ok=True)
    tmp_path = out_dir / f".deck_slide_{page_index + 1:03d}.png"
    pil_image.save(tmp_path, format="PNG")
    return tmp_path


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Export a PDF deck to a PowerPoint (.pptx) with one PDF page per slide (as images)."
    )
    parser.add_argument(
        "--pdf",
        type=Path,
        default=Path(__file__).with_name("deck.pdf"),
        help="Path to input PDF (default: presentation/deck.pdf)",
    )
    parser.add_argument(
        "--pptx",
        type=Path,
        default=Path(__file__).with_name("deck.pptx"),
        help="Path to output PPTX (default: presentation/deck.pptx)",
    )
    parser.add_argument(
        "--dpi",
        type=int,
        default=200,
        help="Render DPI for PDF pages (higher = sharper, larger PPTX). Default: 200",
    )

    args = parser.parse_args()
    export_pdf_to_pptx(args.pdf, args.pptx, dpi=args.dpi)


if __name__ == "__main__":
    main()
